from pptx import Presentation
from pptx.util import Inches

# Create a presentation object
prs = Presentation()

# Slide Titles and Content
slides_content = [
    ("M&A Sell-Side Pitch", "For [Apparel Company Name]"),
    ("Executive Summary", "Key highlights and value proposition."),
    ("Market Overview", "Insights into the current market landscape and trends."),
    ("Company Overview", "Introduction to [Apparel Company Name], history, and mission."),
    ("Financial Overview", "Summary of financial performance and growth."),
    ("Product Portfolio", "Overview of product lines and key offerings."),
    ("Strategic Fit", "How [Apparel Company Name] complements the buyer's portfolio."),
    ("Operational Synergies", "Potential operational efficiencies and synergies."),
    ("Financial Synergies", "Financial benefits of the acquisition."),
    ("Valuation", "Valuation metrics and methodology."),
    ("Next Steps", "Proposed timeline and next steps in the M&A process."),
]

# Generate Slides
for title, content in slides_content:
    slide_layout = prs.slide_layouts[5]  # Using a title and content layout
    slide = prs.slides.add_slide(slide_layout)
    title_placeholder = slide.shapes.title
    body_placeholder = slide.placeholders[1]
    
    title_placeholder.text = title
    body_placeholder.text = content

# Save the presentation
prs.save('ma_sell_side_pitch_apparel_company.pptx')
